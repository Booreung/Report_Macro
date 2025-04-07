from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

code = "BLO"
view_id = "BLO0X00XX"
view_name = "테스트입니다."
ver = "v0.1"


개정번호 = "Ver 0.1"
개정일자 = "2025-04-07"
작성자 = "김승민"

target_columns = {
    0: 개정번호,
    2: 개정일자,
    3: 작성자
}

ppt_dir = r"C:\Users\SMKIM\Desktop\새 폴더"
my_dir = r"C:\Users\SMKIM\Desktop\새 폴더 (2)"

os.makedirs(my_dir, exist_ok=True)

ppt_files = [f for f in os.listdir(ppt_dir) if f.endswith(".pptx")]

if not ppt_files:
    print("처리할 파일이 없습니다.")
    exit()


for file in ppt_files:
    org_dir = os.path.join(ppt_dir,file)

    new_file_name = f"화면설계서{code}-{view_id}-{view_name}-{ver}.pptx"
    save_dir = os.path.join(my_dir,new_file_name)


    prs = Presentation(org_dir)

    target_slides = prs.slides[1]
  
    for shape in target_slides.shapes:
        if shape.has_table:
            table = shape.table
            
            # 새로운 행 삽입 (마지막 행을 복제해 빈 값으로 초기화)
            if len(table.rows) >= 2:
                for col_index, text in target_columns.items():
                    cell = table.cell(1,col_index)
                    cell.text = ""

                    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER 

                    if not paragraph.runs:
                        run = paragraph.add_run()
                    else:
                        run = paragraph.runs[0]

                    run.text = text

                    run.font.name = "맑은 고딕"
                    run.font.size = Pt(10)
                    run.font.bold = False
                    run.font.color.rgb = RGBColor(0,0,0)                   

            else:
                print("⚠️ 표에 쓸 수 있는 행이 부족합니다. PPT에서 빈 행을 더 추가해 주세요.")
            break

prs.save(save_dir)
print(f"표 자동 수정 및 저장 완료! -> 저장경로 = {save_dir}")